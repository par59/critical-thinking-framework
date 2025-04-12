from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import pyttsx3
import time
import os

# Initialize text-to-speech engine
engine = pyttsx3.init()
engine.setProperty('rate', 150)  # Speed of speech
engine.setProperty('volume', 1.0)  # Volume (0.0 to 1.0)

# Create output directories if they don't exist
output_dir = 'mcq_output'
audio_dir = os.path.join(output_dir, 'audio')
os.makedirs(output_dir, exist_ok=True)
os.makedirs(audio_dir, exist_ok=True)

def create_audio(text, filename):
    """Create an audio file from text"""
    audio_path = os.path.join(audio_dir, f'{filename}.mp3')
    engine.save_to_file(text, audio_path)
    engine.runAndWait()

# Create a new presentation
prs = Presentation()

# Define slide layouts
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
bullet_slide_layout = prs.slide_layouts[1]

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Set title font size and color
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Set subtitle font size and color
    subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(51, 51, 51)
    
    # Create audio for title slide
    create_audio(f"{title}. {subtitle}", "title")

def add_question_slide(question_number, question, options):
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    # Set title (question number)
    title_shape.text = f"Question {question_number}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Set content (question and options)
    tf = content_shape.text_frame
    tf.clear()
    
    # Add question
    p = tf.add_paragraph()
    p.text = question
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Add options
    for option in options:
        p = tf.add_paragraph()
        p.text = option
        p.font.size = Pt(20)
        p.level = 1
    
    # Add "Click for Answer" text
    p = tf.add_paragraph()
    p.text = "Click for Answer →"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.font.italic = True
    
    # Create audio for question and options
    audio_text = f"Question {question_number}. {question}"
    for option in options:
        audio_text += f". {option}"
    audio_text += ". Please take 15 seconds to think about your answer."
    create_audio(audio_text, f"question_{question_number}")

def add_answer_slide(question_number, correct_answer):
    slide = prs.slides.add_slide(content_slide_layout)
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    # Set title
    title_shape.text = f"Answer to Question {question_number}"
    title_shape.text_frame.paragraphs[0].font.size = Pt(32)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Set content (correct answer)
    tf = content_shape.text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = f"✔️ {correct_answer}"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 128, 0)
    p.font.bold = True
    
    # Create audio for answer
    create_audio(f"The correct answer is {correct_answer}", f"answer_{question_number}")

# Title slide
add_title_slide("Research Aptitude MCQs", "UGC NET Preparation")

# Add MCQs
mcqs = [
    {
        "question": "Which software is used to check plagiarism?",
        "options": ["a) Sound Forge", "b) Grammarly", "c) Turnitin", "d) Fast Pencil"],
        "correct_answer": "c) Turnitin"
    },
    {
        "question": "What is the term for using someone else's research data without permission?",
        "options": ["a) Citation", "b) Paraphrasing", "c) Plagiarism", "d) Referencing"],
        "correct_answer": "c) Plagiarism"
    },
    {
        "question": "Which research property allows results to apply beyond the specific sample?",
        "options": ["a) Internal Validity", "b) Convergent Validity", "c) Divergent Validity", "d) External Validity"],
        "correct_answer": "d) External Validity"
    },
    {
        "question": "Which is the correct sequence of research steps?",
        "options": [
            "a) Problem → Data Collection → Review → Report",
            "b) Review → Data Collection → Problem → Report",
            "c) Problem → Review → Data Collection → Analysis → Report",
            "d) Problem → Report → Review → Data Collection"
        ],
        "correct_answer": "c) Problem → Review → Data Collection → Analysis → Report"
    },
    {
        "question": "Which of the following is least sensitive to research ethics?",
        "options": ["a) Identifying Variables", "b) Data Collection", "c) Analysis", "d) Reporting"],
        "correct_answer": "a) Identifying Variables"
    },
    {
        "question": "Quantitative research is associated with which approach?",
        "options": ["a) Ethnographic", "b) Unstructured", "c) Structured", "d) Flexible"],
        "correct_answer": "c) Structured"
    },
    {
        "question": "Which research type is used for immediate application?",
        "options": ["a) Conceptual", "b) Action", "c) Fundamental", "d) Theoretical"],
        "correct_answer": "b) Action"
    },
    {
        "question": "Which research paradigm interprets reality from participants' perspectives?",
        "options": ["a) Experimental", "b) Ethnographic", "c) Quantitative", "d) Descriptive"],
        "correct_answer": "b) Ethnographic"
    },
    {
        "question": "What type of research is done by observing students in hostels?",
        "options": ["a) Experimental", "b) Case Study", "c) Participant Observation", "d) Ethnography"],
        "correct_answer": "c) Participant Observation"
    },
    {
        "question": "Which research types bridge theory and practice?",
        "options": ["a) Fundamental and Applied", "b) Applied and Conceptual", "c) Conceptual and Action", "d) Fundamental and Historical"],
        "correct_answer": "a) Fundamental and Applied"
    },
    {
        "question": "Which method is best to study teachers' and students' adjustment patterns?",
        "options": ["a) Experimental", "b) Case Study and Ethnographic", "c) Quantitative", "d) Action Research"],
        "correct_answer": "b) Case Study and Ethnographic"
    },
    {
        "question": "Describing why stress causes heart attacks is an example of:",
        "options": ["a) Fundamental Research", "b) Conceptual Research", "c) Explanatory Research", "d) Descriptive Research"],
        "correct_answer": "c) Explanatory Research"
    },
    {
        "question": "Which method is used when a child's feeding method is studied post anxiety?",
        "options": ["a) Case Study", "b) Experimental", "c) Ex Post Facto", "d) Clinical"],
        "correct_answer": "c) Ex Post Facto"
    },
    {
        "question": "In structured interviews, questions are:",
        "options": ["a) Flexible", "b) Random", "c) Pre-determined", "d) Subjective"],
        "correct_answer": "c) Pre-determined"
    },
    {
        "question": "Chi-square test is used to measure:",
        "options": ["a) Mean comparison", "b) Correlation", "c) Goodness of Fit", "d) Hypothesis"],
        "correct_answer": "c) Goodness of Fit"
    },
    {
        "question": "Which test compares means of two groups?",
        "options": ["a) ANOVA", "b) T-test", "c) Chi-square", "d) Z-test"],
        "correct_answer": "b) T-test"
    },
    {
        "question": "What does ANOVA stand for?",
        "options": ["a) Analysis of Samples", "b) Analysis of Validity", "c) Analysis of Variance", "d) Advanced Number Verification"],
        "correct_answer": "c) Analysis of Variance"
    },
    {
        "question": "Which is a non-parametric test for comparing two means?",
        "options": ["a) T-test", "b) Chi-square", "c) Mann-Whitney U test", "d) ANOVA"],
        "correct_answer": "c) Mann-Whitney U test"
    },
    {
        "question": "Which method compares means of more than two groups?",
        "options": ["a) T-test", "b) G-test", "c) ANOVA", "d) Chi-square"],
        "correct_answer": "c) ANOVA"
    },
    {
        "question": "Which is a non-probability sampling method?",
        "options": ["a) Simple Random", "b) Stratified", "c) Snowball", "d) Cluster"],
        "correct_answer": "c) Snowball"
    },
    {
        "question": "Sampling where every individual has equal chance is:",
        "options": ["a) Cluster Sampling", "b) Stratified Sampling", "c) Systematic Sampling", "d) Simple Random Sampling"],
        "correct_answer": "d) Simple Random Sampling"
    },
    {
        "question": "Choosing a sample due to ease of access is called:",
        "options": ["a) Snowball Sampling", "b) Judgment Sampling", "c) Convenience Sampling", "d) Stratified Sampling"],
        "correct_answer": "c) Convenience Sampling"
    },
    {
        "question": "What does participant observation involve?",
        "options": ["a) Surveys", "b) Experiments", "c) Deep observation by being part of the setting", "d) Watching from a distance"],
        "correct_answer": "c) Deep observation by being part of the setting"
    },
    {
        "question": "Which of these is a qualitative method?",
        "options": ["a) Statistical Tests", "b) Surveys", "c) Ethnography", "d) ANOVA"],
        "correct_answer": "c) Ethnography"
    },
    {
        "question": "What type of sampling divides the population into groups first?",
        "options": ["a) Random Sampling", "b) Cluster Sampling", "c) Stratified Sampling", "d) Convenience Sampling"],
        "correct_answer": "c) Stratified Sampling"
    },
    {
        "question": "Which research starts with a problem already occurred and works backwards?",
        "options": ["a) Experimental", "b) Action", "c) Ex Post Facto", "d) Longitudinal"],
        "correct_answer": "c) Ex Post Facto"
    },
    {
        "question": "What is the first step in the research process?",
        "options": ["a) Data Analysis", "b) Review of Literature", "c) Identifying the Research Problem", "d) Writing Report"],
        "correct_answer": "c) Identifying the Research Problem"
    },
    {
        "question": "What is the aim of reviewing research literature?",
        "options": ["a) Prove hypothesis", "b) Collect data", "c) Generate and validate research questions", "d) Interpret graphs"],
        "correct_answer": "c) Generate and validate research questions"
    },
    {
        "question": "Which of the following is NOT a probability sampling method?",
        "options": ["a) Stratified", "b) Cluster", "c) Snowball", "d) Systematic"],
        "correct_answer": "c) Snowball"
    },
    {
        "question": "In which sampling are groups selected, not individuals?",
        "options": ["a) Simple Random", "b) Cluster Sampling", "c) Convenience Sampling", "d) Stratified Sampling"],
        "correct_answer": "b) Cluster Sampling"
    },
    {
        "question": "Which research uses real-life settings to explore problems deeply?",
        "options": ["a) Laboratory", "b) Experimental", "c) Case Study", "d) Survey"],
        "correct_answer": "c) Case Study"
    },
    {
        "question": "Which of the following is a parametric test?",
        "options": ["a) Chi-square test", "b) Mann-Whitney U test", "c) Kruskal-Wallis test", "d) ANOVA"],
        "correct_answer": "d) ANOVA"
    },
    {
        "question": "Which test is used to compare more than two group means?",
        "options": ["a) T-test", "b) G-test", "c) ANOVA", "d) Z-test"],
        "correct_answer": "c) ANOVA"
    },
    {
        "question": "The t-test is used for:",
        "options": ["a) Two groups' means comparison", "b) Correlation", "c) Proving hypothesis", "d) Frequency analysis"],
        "correct_answer": "a) Two groups' means comparison"
    },
    {
        "question": "Which method is cyclic in nature involving Plan, Act, Observe, Reflect?",
        "options": ["a) Experimental Research", "b) Case Study", "c) Action Research", "d) Historical Research"],
        "correct_answer": "c) Action Research"
    },
    {
        "question": "Which of the following is NOT a step in the research process?",
        "options": ["a) Data Collection", "b) Creating Syllabus", "c) Hypothesis Testing", "d) Reporting Results"],
        "correct_answer": "b) Creating Syllabus"
    },
    {
        "question": "The purpose of a hypothesis in research is to:",
        "options": ["a) Give conclusion", "b) Propose possible explanations", "c) Collect data", "d) Review literature"],
        "correct_answer": "b) Propose possible explanations"
    },
    {
        "question": "Chi-square test is a:",
        "options": ["a) Parametric Test", "b) Non-parametric Test", "c) Field Study", "d) Quantitative method"],
        "correct_answer": "b) Non-parametric Test"
    },
    {
        "question": "The main objective of sampling is to:",
        "options": ["a) Reduce errors", "b) Save time and resources", "c) Avoid analysis", "d) Avoid hypothesis"],
        "correct_answer": "b) Save time and resources"
    },
    {
        "question": "Which sampling method involves referral from one participant to another?",
        "options": ["a) Cluster Sampling", "b) Snowball Sampling", "c) Stratified Sampling", "d) Quota Sampling"],
        "correct_answer": "b) Snowball Sampling"
    },
    {
        "question": "Which method involves selecting only available subjects?",
        "options": ["a) Quota Sampling", "b) Stratified Sampling", "c) Convenience Sampling", "d) Systematic Sampling"],
        "correct_answer": "c) Convenience Sampling"
    },
    {
        "question": "ANOVA test compares:",
        "options": ["a) Medians", "b) Proportions", "c) Variances", "d) Frequencies"],
        "correct_answer": "c) Variances"
    },
    {
        "question": "Ethnographic research belongs to which category?",
        "options": ["a) Quantitative", "b) Numerical", "c) Qualitative", "d) Survey"],
        "correct_answer": "c) Qualitative"
    },
    {
        "question": "Which sampling gives each unit equal and independent chance?",
        "options": ["a) Systematic Sampling", "b) Convenience Sampling", "c) Simple Random Sampling", "d) Quota Sampling"],
        "correct_answer": "c) Simple Random Sampling"
    },
    {
        "question": "In what method are participants observed in their natural setting?",
        "options": ["a) Participant Observation", "b) Simulation", "c) Experimental", "d) Survey"],
        "correct_answer": "a) Participant Observation"
    },
    {
        "question": "Which step comes after hypothesis testing?",
        "options": ["a) Data Collection", "b) Statement of Objectives", "c) Interpretation of Results", "d) Literature Review"],
        "correct_answer": "c) Interpretation of Results"
    },
    {
        "question": "Which of the following is a descriptive method?",
        "options": ["a) Experimental", "b) Survey", "c) Action", "d) Case Control"],
        "correct_answer": "b) Survey"
    },
    {
        "question": "Participant observation is used in:",
        "options": ["a) Laboratory experiments", "b) Natural settings", "c) Mathematical models", "d) Archival research"],
        "correct_answer": "b) Natural settings"
    },
    {
        "question": "The chi-square test checks:",
        "options": ["a) Differences in means", "b) Goodness of fit", "c) Changes over time", "d) Random sampling"],
        "correct_answer": "b) Goodness of fit"
    },
    {
        "question": "Which sampling technique is used when population is divided into groups?",
        "options": ["a) Stratified Sampling", "b) Snowball Sampling", "c) Purposive Sampling", "d) Random Sampling"],
        "correct_answer": "a) Stratified Sampling"
    }
]

# Add all MCQs to slides
for i, mcq in enumerate(mcqs, 1):
    print(f"Creating slides and audio for question {i}...")
    # Add question slide
    add_question_slide(i, mcq["question"], mcq["options"])
    # Add answer slide
    add_answer_slide(i, mcq["correct_answer"])

print("Creating PowerPoint presentation...")
# Save the presentation
pptx_path = os.path.join(output_dir, 'research_aptitude_mcqs.pptx')
prs.save(pptx_path)

print("\nDone! The presentation has been created with audio narration.")
print(f"\nFiles created:")
print(f"1. PowerPoint presentation: {pptx_path}")
print(f"2. Audio files: {audio_dir}/*.mp3")
print("\nTo use the presentation with audio:")
print("1. Open the presentation in PowerPoint")
print("2. For each slide:")
print("   - Add the corresponding audio file from the 'audio' folder")
print("   - Set the audio to play automatically")
print("3. Go to the 'Slide Show' tab")
print("4. Click on 'Set Up Slide Show'")
print("5. Under 'Show type', select 'Browsed at a kiosk (full screen)'")
print("6. Click OK and start the slide show") 