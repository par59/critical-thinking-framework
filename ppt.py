from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a new presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1]
bullet_slide_layout = prs.slide_layouts[1]

def add_title_slide(title, subtitle):
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle

def add_bullet_slide(title, bullets):
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = title
    content = slide.placeholders[1].text_frame
    content.clear()
    for bullet in bullets:
        p = content.add_paragraph()
        p.text = bullet
        p.level = 0

# Slide 1: Title
add_title_slide("National Assessment and Accreditation Council (NAAC)", 
                "UGC NET Education – Key Concepts & Criteria")

# Slide 2: Introduction
add_bullet_slide("Welcome & Overview", [
    "Importance of NAAC in higher education",
    "Relevance for UGC NET exams"
])

# Slide 3: What is NAAC?
add_bullet_slide("What is NAAC?", [
    "Full Form: National Assessment and Accreditation Council",
    "Established: 1994 by UGC, based on NPE 1986",
    "Headquarters: Bangalore"
])

# Slide 4: NAAC Key Authorities
add_bullet_slide("NAAC Key Authorities", [
    "General Council (GC) – Chaired by UGC Chairperson: Prof. M. Jagadesh Kumar",
    "Executive Committee (EC) – Chairperson: Dr. Bhushan Patwardhan"
])

# Slide 5: NAAC Logo Meaning
add_bullet_slide("NAAC Logo – What it Signifies", [
    "Excellence = Quality",
    "Credibility = Reliability",
    "Relevance = Contemporary Utility"
])

# Slide 6: NAAC’s Mission
add_bullet_slide("NAAC’s Mission – 5 Highlights", [
    "Periodic assessment of teaching-learning",
    "Promote research-based education",
    "Innovation & technology use",
    "Student-centered programs",
    "Quality enhancement"
])

# Slide 7: NAAC Accreditation Process
add_bullet_slide("NAAC Accreditation Process", [
    "1. Institution Data Submission",
    "2. Data Validation & Verification",
    "3. ICT & Digital Evidence Check",
    "4. Student Feedback Collection",
    "5. Scoring & Grading"
])

# Slide 8: 7 Assessment Criteria
add_bullet_slide("7 Assessment Criteria", [
    "1. Curricular Aspects",
    "2. Teaching-Learning & Evaluation",
    "3. Research, Innovation & Extension",
    "4. Infrastructure & Learning Resources",
    "5. Student Support & Progression",
    "6. Governance, Leadership & Management",
    "7. Institutional Values & Best Practices"
])

# Slide 9: Criteria Weightage by Institution
add_bullet_slide("Criteria with Highest Weightage", [
    "Universities: Research, Innovation & Extension",
    "Autonomous Institutions: Teaching & Learning",
    "Affiliated Colleges: Teaching & Learning"
])

# Slide 10: Grading System
add_bullet_slide("Grading System (GPA ➜ Grade)", [
    "3.51 - 4.00   ➜ A++",
    "3.26 - 3.50   ➜ A+",
    "3.01 - 3.25   ➜ A",
    "2.76 - 3.00   ➜ B++",
    "2.51 - 2.75   ➜ B+",
    "2.01 - 2.50   ➜ B",
    "≤ 1.50        ➜ Not Accredited (Unsatisfactory)"
])

# Slide 11: Why NAAC Matters
add_bullet_slide("Why NAAC Matters", [
    "Quality assurance",
    "Student trust",
    "Institutional improvement"
])

# Slide 12: Conclusion
add_bullet_slide("Conclusion & Takeaways", [
    "7 Criteria are the core of evaluation",
    "Grades reflect quality education",
    "NAAC encourages continuous improvement"
])

# Save presentation
pptx_file = "/mnt/data/NAAC_UGC_NET_Presentation.pptx"
prs.save(pptx_file)
pptx_file
