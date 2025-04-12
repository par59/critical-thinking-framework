from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation and blank slide
prs = Presentation()
slide_layout = prs.slide_layouts[6]
slide = prs.slides.add_slide(slide_layout)

# Steps and emojis as placeholders
steps = ["Acquire", "Understand", "Doubt", "Transform", "Analyze", "Present"]
icons = ["📥", "🎓", "❓", "⚙️", "🧠", "📊"]

# Style constants
circle_diameter = Inches(1.3)
spacing = Inches(0.7)
start_x = Inches(0.5)
y_position = Inches(2)

# Helper function for text formatting
def format_textbox(shape, text, font_size=14, bold=True, color=RGBColor(255, 255, 255)):
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    font = run.font
    font.size = Pt(font_size)
    font.bold = bold
    font.color.rgb = color

# Function to add shadow to the shapes
def add_shadow(shape):
    shadow = shape.shadow
    shadow.inherit = False
    shadow.blur_radius = Inches(0.15)
    shadow.distance = Inches(0.2)
    shadow.angle = 45
    shadow.fore_color.rgb = RGBColor(0, 0, 0)  # Black shadow

# Add each step
for i, (step, icon) in enumerate(zip(steps, icons)):
    left = start_x + i * (circle_diameter + spacing)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, y_position, circle_diameter, circle_diameter)
    
    # Apply solid color fill and shadow to simulate 3D effect
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(91, 155, 213)  # Blue color for circle
    circle.line.color.rgb = RGBColor(255, 255, 255)  # White border
    circle.line.width = Pt(2)

    add_shadow(circle)  # Add shadow to circle

    icon_box = slide.shapes.add_textbox(left, y_position - Inches(0.6), circle_diameter, Inches(0.5))
    format_textbox(icon_box, icon, font_size=28, bold=False)

    label_box = slide.shapes.add_textbox(left, y_position + circle_diameter, circle_diameter, Inches(0.5))
    format_textbox(label_box, step, font_size=12, bold=True, color=RGBColor(0, 0, 0))

# Add arrow with shadow
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(0.3), Inches(4), Inches(9), Inches(0.4))
arrow.fill.solid()
arrow.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue color for arrow
arrow.line.color.rgb = RGBColor(255, 255, 255)  # White border

add_shadow(arrow)  # Add shadow to arrow

# Save to specified path
prs.save(r"D:\indiannn\lovely\process_flow_slide.pptx")
