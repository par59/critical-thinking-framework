from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_title_slide(prs, title_text, subtitle_text):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = title_text
    subtitle.text = subtitle_text
    return slide

def create_question_slide(prs, question_number, question_text, options, answer):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    
    # Add title
    title = slide.shapes.title
    title.text = f"Question {question_number}"
    
    # Add content
    content = slide.placeholders[1]
    tf = content.text_frame
    
    # Add question
    p = tf.add_paragraph()
    p.text = question_text
    p.font.size = Pt(24)
    
    # Add options
    for option in options:
        p = tf.add_paragraph()
        p.text = option
        p.font.size = Pt(20)
        p.level = 1
    
    # Add answer
    p = tf.add_paragraph()
    p.text = f"Answer: {answer}"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 128, 0)
    
    return slide

def main():
    prs = Presentation()
    
    # Create title slide
    create_title_slide(prs, "Research Aptitude MCQs", "UGC NET Style Questions")
    
    # Questions and answers
    questions = [
        {
            "q": "Which software is used to check plagiarism?",
            "options": ["a) Sound Forge", "b) Grammarly", "c) Turnitin", "d) Fast Pencil"],
            "ans": "c) Turnitin"
        },
        {
            "q": "What is the term for using someone else's research data without permission?",
            "options": ["a) Citation", "b) Paraphrasing", "c) Plagiarism", "d) Referencing"],
            "ans": "c) Plagiarism"
        },
        # Add all 50 questions here...
    ]
    
    # Create slides for each question
    for i, q in enumerate(questions, 1):
        create_question_slide(prs, i, q["q"], q["options"], q["ans"])
    
    # Save the presentation
    prs.save('Research_Aptitude_MCQs.pptx')

if __name__ == "__main__":
    main() 